#!/usr/bin/env python3
"""
Kutumb — Hackathon Presentation
Theme: White · Red · Green  |  Artistic, visual-first
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════════════════════════
# COLOUR TOKENS
# ═══════════════════════════════════════════════════════════════
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)   # deep carmine
RED_SOFT    = RGBColor(0xFD, 0xED, 0xEB)   # blush
RED_DARK    = RGBColor(0x8B, 0x1A, 0x14)
GREEN       = RGBColor(0x27, 0xAE, 0x60)   # emerald
GREEN_SOFT  = RGBColor(0xE8, 0xF8, 0xEE)   # mint
GREEN_DARK  = RGBColor(0x17, 0x6B, 0x3B)
DARK        = RGBColor(0x12, 0x12, 0x1A)   # near-black
CHARCOAL    = RGBColor(0x2D, 0x2D, 0x3A)
GRAY        = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY  = RGBColor(0xF0, 0xF2, 0xF5)

# ═══════════════════════════════════════════════════════════════
# PRESENTATION SETUP
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # truly blank

# ═══════════════════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ═══════════════════════════════════════════════════════════════

def bg(slide, col=WHITE):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = col

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(1.5)):
    shp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp

def oval(slide, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    shp = slide.shapes.add_shape(9, int(l), int(t), int(w), int(h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp

def tb(slide, text, l, t, w, h, size=16, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    """Text box — supports \\n for paragraph breaks."""
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name   = font
    return box

def stat_card(slide, number, label, l, t, w=Inches(2.8), h=Inches(1.5),
              num_col=RED, fill_col=RED_SOFT, border_col=RED):
    w, h = int(w), int(h)
    rect(slide, l, t, w, h, fill=fill_col, line=border_col, lw=Pt(2))
    tb(slide, number, l+Inches(0.1), t+Inches(0.08), w-Inches(0.2), int(h*0.56),
       size=36, bold=True, color=num_col, align=PP_ALIGN.CENTER)
    tb(slide, label, l+Inches(0.1), t+int(h*0.58), w-Inches(0.2), int(h*0.38),
       size=10, color=CHARCOAL, align=PP_ALIGN.CENTER)

def badge(slide, text, l, t, w=Inches(1.8), h=Inches(0.36), fill=GREEN, text_col=WHITE):
    rect(slide, l, t, w, h, fill=fill)
    tb(slide, text, l+Inches(0.05), t+Inches(0.04), w-Inches(0.1), h-Inches(0.08),
       size=10, bold=True, color=text_col, align=PP_ALIGN.CENTER)

def circle_label(slide, text, l, t, size=Inches(0.5), fill=RED, fs=13):
    oval(slide, l, t, size, size, fill=fill)
    tb(slide, text, l, t+Inches(0.04), size, size-Inches(0.08),
       size=fs, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def avatar(slide, letter, l, t, size=Inches(0.65), fill=RED):
    oval(slide, l, t, size, size, fill=fill)
    tb(slide, letter, l, t+Inches(0.08), size, size-Inches(0.1),
       size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def section_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.08), fill=RED)
    tb(slide, title, Inches(0.45), Inches(0.15), Inches(10), Inches(0.85),
       size=30, bold=True, color=DARK)
    if subtitle:
        tb(slide, subtitle, Inches(0.45), Inches(0.92), Inches(11), Inches(0.38),
           size=12, italic=True, color=GRAY)
    rect(slide, Inches(0.45), Inches(1.28), Inches(1.8), Inches(0.04), fill=GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Decorative blobs
oval(s, Inches(9.5), Inches(-1.8), Inches(5.2), Inches(5.2), fill=RED_SOFT)
oval(s, Inches(-1.5), Inches(5.0), Inches(4.0), Inches(4.0), fill=GREEN_SOFT)
oval(s, Inches(7.2), Inches(5.2), Inches(2.2), Inches(2.2), fill=GREEN_SOFT)
# Top stripe
rect(s, 0, 0, prs.slide_width, Inches(0.1), fill=RED)
# Green bottom stripe
rect(s, 0, Inches(7.3), prs.slide_width, Inches(0.2), fill=GREEN)

tb(s, "WITCHHUNT 2026  ·  HEALTH & WELLBEING  ·  HOPEWORKS",
   Inches(0.5), Inches(0.15), Inches(7), Inches(0.38),
   size=9, color=GRAY)

# Hindi title
tb(s, "\u0915\u0941\u091f\u0941\u0902\u092c", Inches(0.45), Inches(0.58), Inches(6), Inches(2.0),
   size=90, bold=True, color=RED, font="Nirmala UI")
tb(s, "K U T U M B", Inches(0.48), Inches(2.52), Inches(6), Inches(0.9),
   size=40, bold=True, color=DARK)

rect(s, Inches(0.48), Inches(3.38), Inches(3.8), Inches(0.06), fill=GREEN)

tb(s, "Ghar se door, dil ke paas.", Inches(0.48), Inches(3.52), Inches(6), Inches(0.45),
   size=17, italic=True, color=GREEN_DARK)
tb(s, "Away from home, close to the heart.", Inches(0.48), Inches(3.95), Inches(6), Inches(0.38),
   size=12, color=GRAY)
tb(s, "AI-powered family health companion\nVoice-first  ·  100+ languages  ·  Privacy-first",
   Inches(0.48), Inches(4.42), Inches(6), Inches(0.75),
   size=13, color=CHARCOAL)

tb(s, "Amrit  ·  Nishita", Inches(0.48), Inches(5.3), Inches(3.5), Inches(0.38),
   size=13, bold=True, color=RED)
tb(s, "HopeWorks Hackathon  ·  April 2026", Inches(0.48), Inches(5.65), Inches(3.5), Inches(0.32),
   size=10, color=GRAY)

stat_card(s, "1.4B", "Indians to reach first",
          Inches(8.1), Inches(1.4), w=Inches(4.5), h=Inches(1.55),
          num_col=RED, fill_col=RED_SOFT, border_col=RED)
stat_card(s, "7,100", "Languages to serve globally",
          Inches(8.1), Inches(3.08), w=Inches(4.5), h=Inches(1.55),
          num_col=GREEN_DARK, fill_col=GREEN_SOFT, border_col=GREEN)
stat_card(s, "\u20b90 / $0", "Health data stored on our servers",
          Inches(8.1), Inches(4.76), w=Inches(4.5), h=Inches(1.55),
          num_col=RED, fill_col=RED_SOFT, border_col=RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — MEET THE TEAM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.2), Inches(-0.6), Inches(3.2), Inches(3.2), fill=GREEN_SOFT)
oval(s, Inches(-0.8), Inches(5.2), Inches(2.8), Inches(2.8), fill=RED_SOFT)

section_bar(s, "MEET THE TEAM", "Two builders. One mission. A billion families — and a world that needs better health tools.")

# Card Amrit
rect(s, Inches(0.45), Inches(1.5), Inches(5.9), Inches(5.65), fill=LIGHT_GRAY, line=RED, lw=Pt(2.5))
oval(s, Inches(2.25), Inches(1.65), Inches(2.05), Inches(2.05), fill=RED)
tb(s, "AL", Inches(2.25), Inches(1.78), Inches(2.05), Inches(1.8),
   size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Amrit Lahari", Inches(0.55), Inches(3.78), Inches(5.7), Inches(0.5),
   size=21, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "Full-Stack Developer & AI Engineer", Inches(0.55), Inches(4.25), Inches(5.7), Inches(0.36),
   size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s, "B.Tech 3\u02b3\u1d48 Year  ·  BITS Pilani", Inches(0.55), Inches(4.58), Inches(5.7), Inches(0.32),
   size=11, color=GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(0.9), Inches(5.0), Inches(5.0), Inches(0.04), fill=GREEN)
tb(s, "FastAPI  ·  Anthropic Claude SDK  ·  Firebase\nReact Native  ·  AES-256-GCM  ·  System Architecture",
   Inches(0.55), Inches(5.1), Inches(5.7), Inches(0.72),
   size=10, color=CHARCOAL, align=PP_ALIGN.CENTER, italic=True)

# Card Nishita
rect(s, Inches(6.6), Inches(1.5), Inches(5.9), Inches(5.65), fill=LIGHT_GRAY, line=GREEN, lw=Pt(2.5))
oval(s, Inches(8.55), Inches(1.65), Inches(2.05), Inches(2.05), fill=GREEN)
tb(s, "NA", Inches(8.55), Inches(1.78), Inches(2.05), Inches(1.8),
   size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Nishita Agarwal", Inches(6.7), Inches(3.78), Inches(5.7), Inches(0.5),
   size=21, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "Product Designer & Frontend Engineer", Inches(6.7), Inches(4.25), Inches(5.7), Inches(0.36),
   size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
tb(s, "B.Tech 3\u02b3\u1d48 Year  ·  BITS Pilani", Inches(6.7), Inches(4.58), Inches(5.7), Inches(0.32),
   size=11, color=GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(7.05), Inches(5.0), Inches(5.0), Inches(0.04), fill=RED)
tb(s, "UI/UX Design  ·  React Native  ·  Voice UX\nMultilingual UX Research  ·  Data Viz  ·  Accessibility",
   Inches(6.7), Inches(5.1), Inches(5.7), Inches(0.72),
   size=10, color=CHARCOAL, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(6.8), Inches(-1.2), Inches(5.5), Inches(5.5), fill=RED_SOFT)
oval(s, Inches(-1.0), Inches(4.5), Inches(3.2), Inches(3.2), fill=GREEN_SOFT)

rect(s, 0, 0, prs.slide_width, Inches(0.1), fill=RED)

tb(s, '"How many of us call home every day,\nask \'how are you?\' and get \'I\'m fine\'',
   Inches(0.5), Inches(0.2), Inches(7.0), Inches(1.7),
   size=22, bold=True, color=DARK)
tb(s, "— even when it\u2019s not true.\u201d",
   Inches(0.5), Inches(1.82), Inches(7.0), Inches(0.55),
   size=22, bold=True, color=RED)

rect(s, Inches(0.5), Inches(2.48), Inches(1.5), Inches(0.05), fill=GREEN)
tb(s, "THE GLOBAL PROBLEM", Inches(0.5), Inches(2.62), Inches(4), Inches(0.32),
   size=10, bold=True, color=GREEN_DARK)

tb(s, "Healthcare systems are reactive, not proactive.\nFamilies engage doctors only during crises.\nDoctors receive zero longitudinal data.",
   Inches(0.5), Inches(3.05), Inches(6.3), Inches(1.15),
   size=14, color=CHARCOAL)

tb(s, "India\u2019s doctor-patient ratio: 1:1456\n— nearly 50% worse than WHO minimum.",
   Inches(0.5), Inches(4.28), Inches(6.3), Inches(0.82),
   size=13, color=CHARCOAL)

stat_card(s, "4.5B",  "People lack essential health access (WHO 2023)",
          Inches(8.4), Inches(0.7), w=Inches(4.5), h=Inches(1.6),
          num_col=RED, fill_col=RED_SOFT, border_col=RED)
stat_card(s, "68%",   "Of global disease burden is preventable",
          Inches(8.4), Inches(2.45), w=Inches(4.5), h=Inches(1.6),
          num_col=GREEN_DARK, fill_col=GREEN_SOFT, border_col=GREEN)
stat_card(s, "$8.3T", "Annual economic loss from poor health outcomes",
          Inches(8.4), Inches(4.2), w=Inches(4.5), h=Inches(1.6),
          num_col=RED, fill_col=RED_SOFT, border_col=RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — HOW THE PROBLEM SHOWS UP
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(10.5), Inches(5.0), Inches(3.2), Inches(3.2), fill=GREEN_SOFT)
section_bar(s, "HOW THE PROBLEM SHOWS UP", "Manifestation for users  ·  Impact on the system")

tb(s, "HOW IT MANIFESTS", Inches(0.45), Inches(1.42), Inches(5), Inches(0.32),
   size=10, bold=True, color=RED)
tb(s, "HOW IT IMPACTS", Inches(7.0), Inches(1.42), Inches(5.5), Inches(0.32),
   size=10, bold=True, color=GREEN_DARK)
rect(s, Inches(6.72), Inches(1.48), Inches(0.04), Inches(5.7), fill=LIGHT_GRAY)

manifests = [
    ("01", RED,       "Silent deterioration \u2014 \u2018I\u2019m fine\u2019 hides months of worsening symptoms doctors never see"),
    ("02", GREEN,     "No longitudinal data at consultations. Doctors start fresh every single visit"),
    ("03", RED,       "Language exclusion: 40%+ of patients cannot describe symptoms in the healthcare system\u2019s language"),
    ("04", GREEN,     "Medicine non-adherence costs $300B/year (US alone) \u2014 125,000 preventable deaths annually"),
    ("05", RED,       "Rural invisibility: 3.5B people in areas with zero specialist access"),
]
for i, (num, col, text) in enumerate(manifests):
    y = Inches(1.75) + i * Inches(1.06)
    circle_label(s, num, Inches(0.45), y + Inches(0.1), size=Inches(0.46), fill=col, fs=12)
    tb(s, text, Inches(1.02), y, Inches(5.55), Inches(0.95),
       size=13, color=CHARCOAL)

impacts = [
    ("$300B",  "Annual cost of medication non-adherence globally",   RED,       RED_SOFT,  RED),
    ("80%",    "Of chronic disease could shift to preventive models", GREEN_DARK,GREEN_SOFT,GREEN),
    ("1:1456", "India\u2019s doctor-patient ratio vs WHO\u2019s 1:1000 target",     RED,       RED_SOFT,  RED),
    ("3.5B",   "People in LMICs with zero specialist access",        GREEN_DARK,GREEN_SOFT,GREEN),
]
for i, (num, label, nc, fc, bc) in enumerate(impacts):
    y = Inches(1.72) + i * Inches(1.16)
    stat_card(s, num, label, Inches(7.0), y, w=Inches(5.9), h=Inches(1.0),
              num_col=nc, fill_col=fc, border_col=bc)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(-0.9), Inches(-0.9), Inches(3.2), Inches(3.2), fill=RED_SOFT)
oval(s, Inches(10.8), Inches(5.4), Inches(3.2), Inches(3.2), fill=GREEN_SOFT)

section_bar(s, "KUTUMB \u2014 THE SOLUTION",
            "A Digital Health Custodian \u2014 we hear you when you don\u2019t mean your \u201cI\u2019m fine\u201ds")

# Left description panel
rect(s, Inches(0.45), Inches(1.48), Inches(4.4), Inches(5.68), fill=LIGHT_GRAY, line=RED, lw=Pt(2))
tb(s, "What We Build", Inches(0.65), Inches(1.62), Inches(4.0), Inches(0.38),
   size=12, bold=True, color=RED)
tb(s, "A voice-first AI\nfamily health companion.", Inches(0.65), Inches(2.05), Inches(4.0), Inches(1.2),
   size=21, bold=True, color=DARK)

features_desc = [
    "\u00b7 Listens in any dialect or language",
    "\u00b7 Monitors symptoms, vitals & medicines daily",
    "\u00b7 Alerts family in real time, anywhere on Earth",
    "\u00b7 Generates doctor-ready clinical summaries",
    "\u00b7 Triggers emergency SOS automatically",
    "\u00b7 Works offline in villages with no internet",
    "\u00b7 Health data stays on your phone only",
]
for i, f in enumerate(features_desc):
    tb(s, f, Inches(0.65), Inches(3.3) + i * Inches(0.46), Inches(4.0), Inches(0.42),
       size=12, color=CHARCOAL)

# Feature grid 2×3
features = [
    ("1", "Voice-First NLP",        "100+ Indian dialects\n+ global language expansion",     RED,        RED_SOFT,  RED),
    ("2", "Family Dashboard",       "Real-time health view\nof every family member",          GREEN_DARK, GREEN_SOFT, GREEN),
    ("3", "Genetic Risk Engine",    "Hereditary disease logic\nauto-adjusts per member",      RED,        RED_SOFT,  RED),
    ("4", "One-Tap SOS",            "Family + hospital + ambulance\nalerted in under 60s",    GREEN_DARK, GREEN_SOFT, GREEN),
    ("5", "Smart Medicine Cabinet", "OCR scan · Interaction check\nDose reminders + refills", RED,        RED_SOFT,  RED),
    ("6", "Doctor Referral PDF",    "Auto-generated longitudinal\nhealth summary for visits", GREEN_DARK, GREEN_SOFT, GREEN),
]
col_xs = [Inches(5.12), Inches(9.35)]
for i, (num, title, desc, nc, fc, bc) in enumerate(features):
    ci = i % 2; ri = i // 2
    l = col_xs[ci]
    t = Inches(1.48) + ri * Inches(1.95)
    w = Inches(4.0); h = Inches(1.82)
    rect(s, l, t, w, h, fill=fc, line=bc, lw=Pt(2))
    circle_label(s, num, l+Inches(0.1), t+Inches(0.12), size=Inches(0.44), fill=bc, fs=12)
    tb(s, title, l+Inches(0.64), t+Inches(0.13), w-Inches(0.74), Inches(0.42),
       size=13, bold=True, color=nc)
    tb(s, desc, l+Inches(0.1), t+Inches(0.62), w-Inches(0.2), Inches(1.1),
       size=11, color=CHARCOAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — WHO WE SERVE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.0), Inches(-0.6), Inches(3.0), Inches(3.0), fill=RED_SOFT)
section_bar(s, "WHO WE SERVE", "Three real humans. One platform.")

personas = [
    ("R", "Ramprasad Sharma",    "Retired teacher, 68  ·  Jaipur  ·  Diabetic",
     [("Diabetic",RED),("Hypertensive",RED),("Tech-hesitant",GREEN),("Children abroad",GREEN)],
     "Medicine reminders in Marwari,\nreassure his children",
     "Says \u2018I\u2019m fine\u2019 every call.\nForgets 3 medicines weekly.",
     "Medicine adherence +70%.\nSOS to children in <60s.",
     RED, RED_SOFT, RED),
    ("P", "Priya Mehta",          "Software engineer, 32  ·  Bengaluru / Chicago",
     [("Global caregiver",GREEN),("Time-starved",GREEN),("Data-driven",RED),("Helpless remotely",RED)],
     "Real-time window into parents\u2019\nhealth from 8,000 miles",
     "Can\u2019t attend doctor visits.\nNo hereditary risk insight.",
     "Live family dashboard.\n3h/week saved on remote care.",
     GREEN_DARK, GREEN_SOFT, GREEN),
    ("S", "Savitribai Patil",     "ASHA worker, 38  ·  Vidarbha  ·  Serves 200+ families",
     [("200+ families",RED),("No specialist",RED),("Offline-critical",GREEN),("Trusted locally",GREEN)],
     "Clinical support in Marathi,\noffline, for rural monitoring",
     "Nearest doctor 30 km away.\nPaper records lost.",
     "200 families on screen.\nWrong referrals cut 50%.",
     RED, RED_SOFT, RED),
]

for i, (letter, name, sub, tags, need, pain, gain, nc, fc, bc) in enumerate(personas):
    l = Inches(0.45) + i * Inches(4.35)
    t = Inches(1.48); cw = Inches(4.1); ch = Inches(5.68)
    rect(s, l, t, cw, ch, fill=fc, line=bc, lw=Pt(2))
    avatar(s, letter, l+Inches(1.72), t+Inches(0.12), size=Inches(0.68), fill=bc)
    tb(s, name, l+Inches(0.1), t+Inches(0.88), cw-Inches(0.2), Inches(0.42),
       size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, sub, l+Inches(0.1), t+Inches(1.28), cw-Inches(0.2), Inches(0.38),
       size=9, color=GRAY, align=PP_ALIGN.CENTER)
    # Tags
    for j, (tag, tc) in enumerate(tags):
        tx = l+Inches(0.12) + (j%2)*Inches(1.96)
        ty = t+Inches(1.75) + (j//2)*Inches(0.42)
        badge(s, tag, tx, ty, w=Inches(1.86), h=Inches(0.34), fill=tc)
    y0 = t+Inches(2.62)
    tb(s, "NEED", l+Inches(0.12), y0, cw-Inches(0.2), Inches(0.28),
       size=9, bold=True, color=GREEN_DARK)
    tb(s, need, l+Inches(0.12), y0+Inches(0.27), cw-Inches(0.2), Inches(0.65),
       size=11, color=CHARCOAL)
    tb(s, "PAIN", l+Inches(0.12), y0+Inches(0.98), cw-Inches(0.2), Inches(0.28),
       size=9, bold=True, color=RED)
    tb(s, pain, l+Inches(0.12), y0+Inches(1.25), cw-Inches(0.2), Inches(0.65),
       size=11, color=CHARCOAL)
    tb(s, "GAIN", l+Inches(0.12), y0+Inches(1.96), cw-Inches(0.2), Inches(0.28),
       size=9, bold=True, color=GREEN_DARK)
    tb(s, gain, l+Inches(0.12), y0+Inches(2.23), cw-Inches(0.2), Inches(0.65),
       size=11, color=CHARCOAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — AI BRAIN 4-LAYER STACK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.2), Inches(5.5), Inches(2.8), Inches(2.8), fill=RED_SOFT)
section_bar(s, "THE AI BRAIN \u2014 4-LAYER INTELLIGENCE STACK",
            "Every layer has one job. Together they save lives.")

layers = [
    ("Layer 1", "Multilingual NLP Voice Engine",
     "Google STT + Bhashini ASR transcribes 100+ Indian dialects. Sentiment detects fatigue, stress, low mood from voice tone \u2014 even when user says \u2018I\u2019m fine\u2019. Works offline at 2G.",
     "Bhashini  \u00b7  Whisper large-v3  \u00b7  On-device TFLite",
     RED, RED_SOFT, RED),
    ("Layer 2", "Claude AI Health Memory System (Core Brain)",
     "claude-sonnet-4-6 powers all health intelligence. User\u2019s full encrypted history (.kutumb AES-256-GCM) sent per-session. Returns AI response + structured memory patches + trigger rules. Zero server-side storage.",
     "claude-sonnet-4-6  \u00b7  3-layer prompt  \u00b7  Local AES-256-GCM",
     GREEN_DARK, GREEN_SOFT, GREEN),
    ("Layer 3", "Predictive Health Intelligence & Trigger Engine",
     "Trend analysis detects declining organ scores, recurring symptoms, medicine non-adherence. Hard trigger rules fire locally offline. SpO2 < 88 \u2192 immediate SOS. Blood glucose > 300 \u2192 emergency.",
     "Custom organ scoring (Heart/Brain/Gut/Lungs)  \u00b7  Firestore triggers",
     RED, RED_SOFT, RED),
    ("Layer 4", "Computer Vision \u2014 Prescriptions & Lab Reports",
     "Google Cloud Vision OCR reads prescription photos and PDF lab reports. Extracts drug names, dosages, biomarker values. Critical lab values trigger push notification + family SOS. Doctor referral PDF auto-generated.",
     "Google Cloud Vision API  \u00b7  LLM extraction  \u00b7  WeasyPrint referral PDF",
     GREEN_DARK, GREEN_SOFT, GREEN),
]

for i, (lname, title, desc, tech, nc, fc, bc) in enumerate(layers):
    t = Inches(1.42) + i * Inches(1.47)
    rect(s, Inches(0.45), t, Inches(9.0), Inches(1.3), fill=fc, line=bc, lw=Pt(2))
    rect(s, Inches(0.45), t, Inches(1.15), Inches(1.3), fill=bc)
    tb(s, lname, Inches(0.48), t+Inches(0.44), Inches(1.1), Inches(0.38),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, Inches(1.7), t+Inches(0.08), Inches(7.6), Inches(0.42),
       size=14, bold=True, color=nc)
    tb(s, desc, Inches(1.7), t+Inches(0.52), Inches(7.6), Inches(0.72),
       size=10, color=CHARCOAL)
    rect(s, Inches(9.62), t, Inches(3.5), Inches(1.3), fill=LIGHT_GRAY, line=bc, lw=Pt(1))
    tb(s, tech, Inches(9.75), t+Inches(0.35), Inches(3.2), Inches(0.6),
       size=10, color=nc, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — LANGUAGE SCALABILITY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(9.2), Inches(-1.2), Inches(5.0), Inches(5.0), fill=GREEN_SOFT)
oval(s, Inches(-1.2), Inches(4.8), Inches(3.2), Inches(3.2), fill=RED_SOFT)

section_bar(s, "LANGUAGE SCALABILITY \u2014 7,100 LANGUAGES",
            "Health advice that speaks your language, not just the hospital\u2019s.")

# Giant number
tb(s, "7,100", Inches(0.45), Inches(1.35), Inches(4.5), Inches(1.4),
   size=76, bold=True, color=RED)
tb(s, "languages", Inches(0.48), Inches(2.65), Inches(4.5), Inches(0.42),
   size=18, bold=True, color=GREEN_DARK)

# Quote
rect(s, Inches(4.6), Inches(1.35), Inches(8.5), Inches(1.22), fill=RED_SOFT, line=RED, lw=Pt(1.5))
tb(s, "\u201cThe same AI that asks \u2018Did you take your BP medicine?\u2019 in Marwari\nwill ask the same question in Swahili, Hausa, Tagalog, and Haitian Creole\u2014\ntrained on the same architecture, deployed on the same app.\u201d",
   Inches(4.78), Inches(1.42), Inches(8.15), Inches(1.08),
   size=11, italic=True, color=RED_DARK)

tiers = [
    ("TIER 1","Day 1",       "11 langs  \u00b7  1.2B speakers",
     "Hindi  \u00b7  Marathi  \u00b7  Bengali  \u00b7  Tamil  \u00b7  Telugu  \u00b7  Kannada  \u00b7  Gujarati  \u00b7  Punjabi  \u00b7  Malayalam  \u00b7  Odia  \u00b7  English",
     RED, RED_SOFT),
    ("TIER 2","3\u20136 months","10 langs  \u00b7  800M+ speakers",
     "Swahili  \u00b7  Hausa  \u00b7  Amharic  \u00b7  Yoruba  \u00b7  Tagalog  \u00b7  Bahasa  \u00b7  Vietnamese  \u00b7  Burmese  \u00b7  Nepali  \u00b7  Sinhala",
     GREEN_DARK, GREEN_SOFT),
    ("TIER 3","12\u201324 months","14 lang families  \u00b7  3B+ speakers",
     "Spanish  \u00b7  Portuguese  \u00b7  Arabic  \u00b7  French (Africa)  \u00b7  Mandarin  \u00b7  Urdu  \u00b7  Pashto  \u00b7  Dari  \u00b7  Haitian Creole + dialects",
     RED, RED_SOFT),
    ("TIER 4","36 months+",  "Full global coverage",
     "All 7,100 UNESCO-recognised languages via low-resource NLP, code-switching models, federated on-device fine-tuning",
     GREEN_DARK, GREEN_SOFT),
]
for i, (tier, timeline, speakers, langs, nc, fc) in enumerate(tiers):
    t = Inches(2.88) + i * Inches(1.1)
    rect(s, Inches(0.45), t, Inches(12.5), Inches(0.98), fill=fc, line=nc, lw=Pt(2))
    rect(s, Inches(0.45), t, Inches(1.45), Inches(0.98), fill=nc)
    tb(s, tier,     Inches(0.5),  t+Inches(0.08), Inches(1.35), Inches(0.36), size=10, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, timeline, Inches(0.5),  t+Inches(0.44), Inches(1.35), Inches(0.35), size=10, italic=True,color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, speakers, Inches(2.0),  t+Inches(0.06), Inches(2.5),  Inches(0.36), size=10, bold=True,  color=nc)
    tb(s, langs,    Inches(2.0),  t+Inches(0.44), Inches(10.8), Inches(0.46), size=10, color=CHARCOAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.2), Inches(5.4), Inches(2.8), Inches(2.8), fill=GREEN_SOFT)
section_bar(s, "SYSTEM ARCHITECTURE",
            "Privacy-first by design \u2014 health data never leaves the device")

# Tech stack row
tech_stack = [
    ("React Native", RED), ("FastAPI (Python)", GREEN_DARK), ("claude-sonnet-4-6", RED),
    ("Firebase", GREEN_DARK), ("Google Vision", RED), ("AES-256-GCM", GREEN_DARK),
]
for i, (tech, col) in enumerate(tech_stack):
    badge(s, tech, Inches(0.45) + i*Inches(2.17), Inches(1.38), w=Inches(2.05), h=Inches(0.36), fill=col)

# Architecture flow
arch = [
    ("1","User Device",     "React Native\nEncrypted .kutumb",  RED_SOFT,   RED),
    ("2","FastAPI Backend", "12 modules\n30+ endpoints",        GREEN_SOFT, GREEN),
    ("3","Claude AI",       "claude-sonnet-4-6\nMemory + Patches", RED_SOFT, RED),
    ("4","Firebase",        "Firestore + Storage\nFamily & labs", GREEN_SOFT, GREEN),
    ("5","SOS & Alerts",    "FCM push + SMS\nEmergency relay",  RED_SOFT,   RED),
]
bw = Inches(2.28); bh = Inches(1.55)
for i, (num, title, desc, fc, bc) in enumerate(arch):
    l = Inches(0.3) + i*Inches(2.62); t = Inches(1.92)
    rect(s, l, t, bw, bh, fill=fc, line=bc, lw=Pt(2))
    circle_label(s, num, l+int(bw*0.5)-Inches(0.22), t+Inches(0.08), size=Inches(0.42), fill=bc, fs=12)
    tb(s, title, l+Inches(0.06), t+Inches(0.55), bw-Inches(0.12), Inches(0.38),
       size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, desc, l+Inches(0.06), t+Inches(0.95), bw-Inches(0.12), Inches(0.52),
       size=10, color=CHARCOAL, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, "\u2192", Inches(0.3)+(i+1)*Inches(2.62)-Inches(0.38), t+Inches(0.55),
           Inches(0.38), Inches(0.38), size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Privacy box
rect(s, Inches(0.45), Inches(3.66), Inches(12.5), Inches(0.82), fill=LIGHT_GRAY, line=GREEN, lw=Pt(1.5))
tb(s, "[LOCK]  LOCAL HEALTH MEMORY: .kutumb file (AES-256-GCM) lives ONLY on user\u2019s device. Sent per-session over TLS 1.3. Purged from server memory immediately. DPDP \u00b7 HIPAA \u00b7 GDPR compliant by architecture.",
   Inches(0.65), Inches(3.74), Inches(12.1), Inches(0.66),
   size=11, color=DARK, italic=True)

# Compliance
compliance = [("  DPDP 2023 (India)",GREEN),("  HIPAA (United States)",RED),
              ("  GDPR (European Union)",GREEN),("  Offline-First (2G)",RED)]
for i, (label, col) in enumerate(compliance):
    badge(s, "\u2713"+label, Inches(0.45)+i*Inches(3.25), Inches(4.65), w=Inches(3.0), h=Inches(0.38), fill=col)

tb(s, "12 BACKEND MODULES:", Inches(0.45), Inches(5.25), Inches(3), Inches(0.32),
   size=10, bold=True, color=DARK)
modules = ["Auth","Users","Family","Check-ins","Medicines","Health Scores",
           "Emergency SOS","Referrals PDF","Wearable Sync","Lab OCR","AI Insights","AI Session"]
for i, mod in enumerate(modules):
    col_c = RED if i%2==0 else GREEN_DARK
    badge(s, mod, Inches(0.45)+(i%6)*Inches(2.17), Inches(5.65)+(i//6)*Inches(0.44),
          w=Inches(2.05), h=Inches(0.36), fill=col_c)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — TECHNICAL REQUEST
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.5), Inches(-0.6), Inches(2.6), Inches(2.6), fill=GREEN_SOFT)
section_bar(s, "TECHNICAL REQUEST \u2014 FACILITATED DATASETS",
            "What we have  \u00b7  What we need  \u00b7  How it unlocks global reach")

tb(s, "WHAT WE HAVE (Open Source, Already Live)", Inches(0.45), Inches(1.44), Inches(5.8), Inches(0.34),
   size=10, bold=True, color=GREEN_DARK)

have = [
    ("Bhashini Open API",     "Hindi / Marathi / Bengali STT live in prototype"),
    ("SNOMED CT (free tier)", "Clinical terminology, symptom tagging"),
    ("RxNorm (NLM open)",     "Drug interaction baseline, name normalisation"),
    ("ICMR OGD Portal",       "40 district health statistics for score calibration"),
    ("Whisper open weights",  "Foundation model for all non-Bhashini languages"),
]
for i, (title, desc) in enumerate(have):
    t = Inches(1.88) + i*Inches(1.0)
    rect(s, Inches(0.45), t, Inches(5.9), Inches(0.85), fill=GREEN_SOFT, line=GREEN, lw=Pt(1.5))
    oval(s, Inches(0.56), t+Inches(0.22), Inches(0.38), Inches(0.38), fill=GREEN)
    tb(s, "\u2713", Inches(0.58), t+Inches(0.22), Inches(0.36), Inches(0.36),
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, Inches(1.06), t+Inches(0.06), Inches(5.15), Inches(0.36),
       size=13, bold=True, color=GREEN_DARK)
    tb(s, desc, Inches(1.06), t+Inches(0.44), Inches(5.15), Inches(0.35),
       size=10, color=CHARCOAL)

tb(s, "WHAT WE NEED (Facilitated / AI4India)", Inches(6.65), Inches(1.44), Inches(6.2), Inches(0.34),
   size=10, bold=True, color=RED)

need = [
    ("01","Bhashini ASR Training Corpus",
     "500h/language  \u00b7  12 languages  \u00b7  Speaker-level granularity (age, region, gender)"),
    ("02","Anonymised Longitudinal Health Records",
     "5 years  \u00b7  100K+ patients  \u00b7  Urban + rural split  \u00b7  ICD-10 diagnosis codes"),
    ("03","ICMR Lab Report Template Library",
     "50+ formats  \u00b7  10,000+ anonymised samples  \u00b7  Lab name, biomarker, value, reference range"),
    ("04","Low-Resource Language Audio (LMICs)",
     "100h/language: Swahili, Hausa, Tagalog, Bahasa  \u00b7  Health-domain sentences preferred"),
]
for i, (num, title, desc) in enumerate(need):
    t = Inches(1.88) + i*Inches(1.3)
    rect(s, Inches(6.65), t, Inches(6.42), Inches(1.12), fill=RED_SOFT, line=RED, lw=Pt(1.5))
    circle_label(s, num, Inches(6.75), t+Inches(0.3), size=Inches(0.45), fill=RED, fs=12)
    tb(s, title, Inches(7.32), t+Inches(0.08), Inches(5.65), Inches(0.38),
       size=13, bold=True, color=RED)
    tb(s, desc, Inches(7.32), t+Inches(0.5), Inches(5.65), Inches(0.55),
       size=10, color=CHARCOAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — ENVIRONMENT & SECURITY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(-0.8), Inches(4.8), Inches(2.8), Inches(2.8), fill=GREEN_SOFT)
section_bar(s, "ENVIRONMENT & SECURITY REQUIREMENTS",
            "Access methods  \u00b7  Compliance  \u00b7  Privacy by architecture")

tb(s, "ACCESS METHOD & DATA DELIVERY", Inches(0.45), Inches(1.42), Inches(6), Inches(0.32),
   size=10, bold=True, color=RED)

access = [
    ("A","Secure API (Preferred)",
     "HTTPS REST with OAuth2 / API-key. Data streamed, not cached. All responses via encrypted channel.",
     RED, RED_SOFT, RED),
    ("B","AIKosh Sandbox",
     "Health records processed entirely inside the AIKosh secure sandbox. No data extraction. Results-only output.",
     GREEN_DARK, GREEN_SOFT, GREEN),
    ("C","Encrypted Download",
     "Static corpora as AES-256 ZIP to isolated research environment. Destruction certification on completion.",
     RED, RED_SOFT, RED),
]
for i, (letter, title, desc, nc, fc, bc) in enumerate(access):
    l = Inches(0.45)+i*Inches(4.35); t = Inches(1.78)
    w = Inches(4.1); h = Inches(1.42)
    rect(s, l, t, w, h, fill=fc, line=bc, lw=Pt(2))
    avatar(s, letter, l+Inches(0.12), t+Inches(0.12), size=Inches(0.6), fill=bc)
    tb(s, title, l+Inches(0.85), t+Inches(0.12), w-Inches(0.95), Inches(0.38),
       size=13, bold=True, color=nc)
    tb(s, desc, l+Inches(0.12), t+Inches(0.65), w-Inches(0.22), Inches(0.72),
       size=10, color=CHARCOAL)

tb(s, "COMPLIANCE ARCHITECTURE", Inches(0.45), Inches(3.32), Inches(6), Inches(0.32),
   size=10, bold=True, color=GREEN_DARK)

comp_rows = [
    ("\u2713 DPDP Act 2023 (India)", RED,
     "No health data stored server-side between sessions. Consent gate at onboarding. Right to erasure: delete account = wipe local .kutumb file."),
    ("\u2713 HIPAA (United States)", GREEN_DARK,
     "Zero PHI stored server-side. TLS 1.3 for all data in transit. Zero logging of health content. Business Associate Agreement available for US."),
    ("\u2713 GDPR (European Union)", RED,
     "Explicit opt-in consent. Data minimisation: only sha256 user hash sent to server. Right to portability (export .kutumb). Right to erasure."),
]
for i, (label, nc, desc) in enumerate(comp_rows):
    t = Inches(3.72) + i*Inches(1.12)
    rect(s, Inches(0.45), t, Inches(12.5), Inches(0.98), fill=LIGHT_GRAY, line=nc, lw=Pt(2))
    rect(s, Inches(0.45), t, Inches(2.55), Inches(0.98), fill=nc)
    tb(s, label, Inches(0.55), t+Inches(0.29), Inches(2.35), Inches(0.38),
       size=11, bold=True, color=WHITE)
    tb(s, desc, Inches(3.15), t+Inches(0.14), Inches(9.6), Inches(0.72),
       size=11, color=CHARCOAL)

rect(s, 0, Inches(7.05), prs.slide_width, Inches(0.38), fill=DARK)
tb(s, "Security: AES-256-GCM on device  \u00b7  TLS 1.3 in transit  \u00b7  sha256 user hash (no reversible PII)  \u00b7  Biometric key wrap (iOS Secure Enclave / Android Keystore)",
   Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.3),
   size=10, color=WHITE, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.5), Inches(-0.6), Inches(2.6), Inches(2.6), fill=RED_SOFT)
oval(s, Inches(-0.8), Inches(5.4), Inches(2.4), Inches(2.4), fill=GREEN_SOFT)
section_bar(s, "BUSINESS MODEL, PARTNERSHIPS & RETENTION",
            "Three tiers. Recurring revenue. Global partnership runway.")

pricing = [
    ("FREE",     "\u20b90 / $0",        GREEN,   GREEN_SOFT, GREEN,
     ["Daily voice check-in","Family dashboard (3 members)","Basic medicine reminder","One-tap SOS","Kutumb AI (5/day)"]),
    ("PLUS",     "\u20b9149 / $2.99 mo", RED,     RED_SOFT,   RED,
     ["Unlimited family members","Full AI health memory","Doctor referral PDFs","Lab report OCR","Wearable sync","Priority ambulance SOS"]),
    ("CARE PRO", "\u20b9499 / $7.99 mo", DARK,    LIGHT_GRAY, CHARCOAL,
     ["ASHA/CHW multi-family view","Clinical decision support","Govt. scheme integration","Telemedicine booking","White-label (NHS \u00b7 CGHS \u00b7 Medicaid)"]),
]
for i, (name, price, hc, fc, bc, feats) in enumerate(pricing):
    l = Inches(0.45)+i*Inches(4.35); w = Inches(4.1)
    rect(s, l, Inches(1.48), w, Inches(0.85), fill=hc)
    tb(s, name,  l, Inches(1.5),  w, Inches(0.42), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, price, l, Inches(1.9),  w, Inches(0.38), size=12, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, l, Inches(2.33), w, Inches(2.85), fill=fc, line=bc, lw=Pt(1.5))
    for j, feat in enumerate(feats):
        tb(s, "\u2713  "+feat, l+Inches(0.15), Inches(2.43)+j*Inches(0.44), w-Inches(0.3), Inches(0.4),
           size=12, color=DARK)

rect(s, Inches(0.45), Inches(5.35), Inches(12.5), Inches(0.36), fill=LIGHT_GRAY)
tb(s, "RETENTION & PARTNERSHIPS", Inches(0.65), Inches(5.38), Inches(4), Inches(0.3),
   size=10, bold=True, color=DARK)

partners = [
    ("Pharmeasy / 1mg",       "Medicine refill + streak coupons",         GREEN),
    ("Thyrocare / Redcliffe", "Lab booking from referral PDF",             RED),
    ("WHO / UNICEF",          "UHC 2030 alignment, LMIC rollout",          GREEN),
    ("NHS \u00b7 CGHS \u00b7 Medicaid", "White-label government deployments", RED),
]
for i, (partner, desc, col) in enumerate(partners):
    l = Inches(0.45)+i*Inches(3.25)
    oval(s, l, Inches(5.88), Inches(0.5), Inches(0.5), fill=col)
    tb(s, partner, l+Inches(0.62), Inches(5.88), Inches(2.55), Inches(0.38),
       size=12, bold=True, color=col)
    tb(s, desc, l+Inches(0.62), Inches(6.24), Inches(2.55), Inches(0.32),
       size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — CHALLENGES & RISKS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.2), Inches(5.4), Inches(2.8), Inches(2.8), fill=RED_SOFT)
section_bar(s, "CHALLENGES & RISKS \u2014 WITH MITIGATIONS",
            "We\u2019ve thought hard about what could go wrong.")

challenges = [
    ("01","Medical Ethics & Non-Diagnosis",
     "Users treating AI advice as clinical diagnosis \u2192 delayed treatment.",
     "Hard-coded Non-Doctor Mode. 48h+ persistent symptom forces \u2018Book a Doctor\u2019. Referral PDF language is always \u2018Pointers for your doctor\u2019, never a diagnosis.",
     RED),
    ("02","Global Privacy & Compliance",
     "Health data exposure across jurisdictions: DPDP, HIPAA, GDPR.",
     "Zero health data on servers. AES-256-GCM on device. sha256 user hash only (non-reversible). Compliant by architecture, not just policy.",
     GREEN_DARK),
    ("03","AI Hallucination in Drug Safety",
     "Fabricated drug interactions directly harm patients on critical medications.",
     "Hard contraindication rules fire BEFORE Claude call. Claude instructed: never invent pharmacology. Severity=\u2018severe\u2019 always sets must_consult_doctor=True.",
     RED),
    ("04","Rural Connectivity & Low-End Devices",
     "Core users: 2G networks, \u20b93,000\u2013\u20b96,000 Android phones, no data plan.",
     "Offline-first: trigger rules run locally. Cached STT models. WhatsApp bot fallback, no app install. Core bundle <5MB. Voice-only mode.",
     GREEN_DARK),
]
for i, (num, title, risk, mitigation, col) in enumerate(challenges):
    ri = i//2; ci = i%2
    l = Inches(0.45)+ci*Inches(6.55); t = Inches(1.48)+ri*Inches(2.82)
    w = Inches(6.2); h = Inches(2.62)
    fc = RED_SOFT if col==RED else GREEN_SOFT
    rect(s, l, t, w, h, fill=fc, line=col, lw=Pt(2))
    circle_label(s, num, l+Inches(0.1), t+Inches(0.1), size=Inches(0.48), fill=col, fs=12)
    tb(s, title, l+Inches(0.7), t+Inches(0.1), w-Inches(0.8), Inches(0.4),
       size=14, bold=True, color=col)
    tb(s, "! Risk: "+risk, l+Inches(0.15), t+Inches(0.58), w-Inches(0.25), Inches(0.72),
       size=11, color=CHARCOAL)
    rect(s, l+Inches(0.15), t+Inches(1.35), w-Inches(0.25), Inches(0.04), fill=col)
    tb(s, "\u2713 Mitigation: "+mitigation, l+Inches(0.15), t+Inches(1.46), w-Inches(0.25), Inches(1.08),
       size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — IMPACT PROJECTIONS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.5), Inches(-0.6), Inches(2.6), Inches(2.6), fill=GREEN_SOFT)
section_bar(s, "ANTICIPATED RESULTS & IMPACT PROJECTIONS",
            "What Kutumb achieves \u2014 quantified, per persona, scalable.")

metrics = [
    ("40%",   "Reduction in missed\nmedicine doses\u2014Year 1",           GREEN_DARK, GREEN_SOFT, GREEN),
    ("3\u00d7",    "Doctor visits prepared\nwith longitudinal health data",  RED,        RED_SOFT,   RED),
    ("60%",   "Fewer emergency\nhospitalisations\n(monitored users)",       GREEN_DARK, GREEN_SOFT, GREEN),
    ("$300B", "Global annual opportunity\nin preventable\nmedication waste", RED,        RED_SOFT,   RED),
]
for i, (num, label, nc, fc, bc) in enumerate(metrics):
    stat_card(s, num, label, Inches(0.45)+i*Inches(3.25), Inches(1.45),
              w=Inches(3.05), h=Inches(1.72), num_col=nc, fill_col=fc, border_col=bc)

tb(s, "PERSONA-WISE QUANTIFIED ADVANTAGE", Inches(0.45), Inches(3.38), Inches(6), Inches(0.32),
   size=10, bold=True, color=DARK)

p_rows = [
    ("R","Ramprasad (Elderly Parent)",
     "Medicine adherence +70%. SOS alerts children in <60s. Doctor visits arrive with 6-month health data. Est. \u20b918,000/year saved on avoidable consultations.",
     "5 min/day\nvoice check-in", RED),
    ("P","Priya (Global Caregiver)",
     "Real-time parent health dashboard from anywhere. 3h/week saved on remote care. 50% fewer \u2018emergency flight\u2019 scenarios. Peace of mind, quantified.",
     "2 min/day\ndashboard review", GREEN_DARK),
    ("S","Savitribai (ASHA Worker)",
     "200 families managed from one screen. Wrong referrals cut by 50%. Maternal monitoring enabled. Fully offline-capable, works on 2G feature phone.",
     "Weekly data\nsync in town", RED),
]
for i, (letter, name, impact, habit, col) in enumerate(p_rows):
    t = Inches(3.82) + i*Inches(1.08)
    fc = GREEN_SOFT if col==GREEN_DARK else RED_SOFT
    rect(s, Inches(0.45), t, Inches(12.5), Inches(0.95), fill=fc, line=col, lw=Pt(2))
    avatar(s, letter, Inches(0.55), t+Inches(0.12), size=Inches(0.68), fill=col)
    tb(s, name, Inches(1.35), t+Inches(0.08), Inches(2.55), Inches(0.36),
       size=13, bold=True, color=col)
    tb(s, impact, Inches(4.0), t+Inches(0.08), Inches(7.3), Inches(0.78),
       size=11, color=CHARCOAL)
    rect(s, Inches(11.42), t, Inches(0.04), Inches(0.95), fill=col)
    tb(s, habit, Inches(11.55), t+Inches(0.18), Inches(1.38), Inches(0.58),
       size=9, color=col, italic=True)

rect(s, Inches(0.45), Inches(7.1), Inches(12.5), Inches(0.32), fill=LIGHT_GRAY)
tb(s, "Scalability: Serverless FastAPI (Cloud Run)  \u00b7  Firestore to 1M+ concurrent users  \u00b7  Local AI cuts server costs 80%  \u00b7  Phase 2: SE Asia, Africa, Latin America",
   Inches(0.65), Inches(7.14), Inches(12.1), Inches(0.26), size=9, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — CREDIBILITY & GLOBAL VISION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

oval(s, Inches(11.2), Inches(5.4), Inches(2.8), Inches(2.8), fill=RED_SOFT)
section_bar(s, "WHY KUTUMB: CREDIBILITY & GLOBAL VISION",
            "Built in India. Designed for the world. Trusted by architecture.")

tb(s, "WHY WE\u2019RE CREDIBLE", Inches(0.45), Inches(1.44), Inches(5), Inches(0.32),
   size=10, bold=True, color=RED)

cred = [
    ("A","Complete architecture built",
     "12 FastAPI modules, 30+ endpoints, full Firestore schema, 3-layer Claude prompt. All documented, all buildable now.",
     RED),
    ("B","Privacy-first by engineering",
     "AES-256-GCM local. Zero health data server storage. Meets DPDP, HIPAA, GDPR simultaneously \u2014 by architecture.",
     GREEN_DARK),
    ("C","Claude AI as core brain",
     "Anthropic claude-sonnet-4-6 with hard ethical guardrails \u2014 no diagnosis, must_consult_doctor enforced, never invents pharmacology.",
     RED),
    ("D","Evidence-based projections",
     "Every number grounded in WHO, ICMR, NLM, and peer-reviewed health economics. Not guesses \u2014 cited research.",
     GREEN_DARK),
]
for i, (letter, title, desc, col) in enumerate(cred):
    t = Inches(1.85) + i*Inches(1.3)
    fc = RED_SOFT if col==RED else GREEN_SOFT
    rect(s, Inches(0.45), t, Inches(5.95), Inches(1.12), fill=fc, line=col, lw=Pt(2))
    avatar(s, letter, Inches(0.56), t+Inches(0.22), size=Inches(0.6), fill=col)
    tb(s, title, Inches(1.28), t+Inches(0.08), Inches(4.98), Inches(0.38),
       size=14, bold=True, color=col)
    tb(s, desc, Inches(1.28), t+Inches(0.5), Inches(4.98), Inches(0.58),
       size=11, color=CHARCOAL)

tb(s, "GLOBAL EXPANSION ROADMAP", Inches(6.7), Inches(1.44), Inches(6), Inches(0.32),
   size=10, bold=True, color=GREEN_DARK)

phases = [
    ("Phase 1 (Now)",  "India",              "11 Indic languages, 300M+ families, ASHA + consumer app", RED),
    ("Phase 2 (6mo)",  "SE Asia + E. Africa","Swahili, Hausa, Tagalog, Bahasa. NGO + government partnerships", GREEN_DARK),
    ("Phase 3 (18mo)", "Latin America + MENA","Spanish, Portuguese, Arabic. Telehealth integration. 3B speakers", RED),
    ("Phase 4 (36mo+)","Global",             "7,100 languages. White-label NHS, Medicaid, Ayushman Bharat", GREEN_DARK),
]
for i, (phase, region, detail, col) in enumerate(phases):
    t = Inches(1.85) + i*Inches(1.3)
    fc = RED_SOFT if col==RED else GREEN_SOFT
    rect(s, Inches(6.7), t, Inches(6.28), Inches(1.12), fill=fc, line=col, lw=Pt(2))
    rect(s, Inches(6.7), t, Inches(1.88), Inches(1.12), fill=col)
    tb(s, phase,  Inches(6.75), t+Inches(0.08), Inches(1.78), Inches(0.38),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, region, Inches(6.75), t+Inches(0.48), Inches(1.78), Inches(0.38),
       size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, detail, Inches(8.72), t+Inches(0.2), Inches(4.12), Inches(0.72),
       size=12, color=CHARCOAL)

rect(s, Inches(0.45), Inches(7.1), Inches(12.5), Inches(0.32), fill=GREEN_SOFT, line=GREEN, lw=Pt(1))
tb(s, "HopeWorks alignment: Kutumb reduces catastrophic health expenses \u2014 the single biggest driver of poverty in LMICs. Health equity = economic empowerment.",
   Inches(0.65), Inches(7.14), Inches(12.1), Inches(0.24), size=9, color=GREEN_DARK, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Artistic blobs
oval(s, Inches(9.0), Inches(-2.0), Inches(5.8), Inches(5.8), fill=RED_SOFT)
oval(s, Inches(-2.0), Inches(4.5), Inches(5.0), Inches(5.0), fill=GREEN_SOFT)
oval(s, Inches(5.2), Inches(5.6), Inches(3.2), Inches(3.2), fill=RED_SOFT)

rect(s, 0, 0, prs.slide_width, Inches(0.1), fill=RED)
rect(s, 0, Inches(7.3), prs.slide_width, Inches(0.2), fill=GREEN)

tb(s, "THANK YOU", Inches(0.5), Inches(0.25), Inches(9), Inches(2.0),
   size=82, bold=True, color=DARK)
tb(s, "for your time and attention", Inches(0.5), Inches(2.22), Inches(6), Inches(0.48),
   size=15, italic=True, color=GRAY)

rect(s, Inches(0.5), Inches(2.78), Inches(2.4), Inches(0.07), fill=RED)
rect(s, Inches(3.05), Inches(2.78), Inches(2.4), Inches(0.07), fill=GREEN)

tb(s, "\u0915\u0941\u091f\u0941\u0902\u092c  \u00b7  KUTUMB", Inches(0.5), Inches(2.95), Inches(8), Inches(0.95),
   size=36, bold=True, color=RED, font="Nirmala UI")

taglines = [
    "Ghar se door, dil ke paas.",
    "Away from home, close to the heart.",
    "Built for India. Designed for the world.",
]
cols_t = [GREEN_DARK, GREEN_DARK, GRAY]
for i, (line, col) in enumerate(zip(taglines, cols_t)):
    tb(s, line, Inches(0.5), Inches(3.9)+i*Inches(0.42), Inches(7.5), Inches(0.38),
       size=14, italic=True, color=col)

rect(s, Inches(0.5), Inches(5.22), Inches(4.3), Inches(1.15), fill=LIGHT_GRAY, line=RED, lw=Pt(2))
tb(s, "Team: Amrit & Nishita", Inches(0.65), Inches(5.3), Inches(4.0), Inches(0.42),
   size=16, bold=True, color=DARK)
tb(s, "WitchHunt 2026  \u00b7  Health & Well-being  \u00b7  HopeWorks",
   Inches(0.65), Inches(5.7), Inches(4.0), Inches(0.32), size=11, color=GRAY)

rect(s, Inches(5.1), Inches(5.22), Inches(4.3), Inches(1.15), fill=LIGHT_GRAY, line=GREEN, lw=Pt(2))
tb(s, "Mission: Health for Every Human", Inches(5.25), Inches(5.3), Inches(4.0), Inches(0.42),
   size=14, bold=True, color=GREEN_DARK)
tb(s, "India \u2192 SE Asia \u2192 Africa \u2192 Latin America \u2192 Global\n7,100 languages. 1 platform. 1 goal.",
   Inches(5.25), Inches(5.7), Inches(4.0), Inches(0.62), size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\amrit\hackathons\witchhunt\Kutumb_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
